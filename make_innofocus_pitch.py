"""
Innofocus Photonics Technology — Corporate Financing Pitch Deck
Based on content from https://innofocus.com.au/
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x1A, 0x2E)
BLUE   = RGBColor(0x00, 0x6E, 0xC8)
TEAL   = RGBColor(0x00, 0xB4, 0xD8)
GOLD   = RGBColor(0xF7, 0xA8, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xC0, 0xD0, 0xE0)
DGREY  = RGBColor(0x12, 0x28, 0x3E)
GREEN  = RGBColor(0x2E, 0xCC, 0x71)
RED    = RGBColor(0xE7, 0x4C, 0x3C)

W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=NAVY, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def bg(slide):    rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
def bar(slide):   rect(slide, 0, 0, 0.20, 7.5, fill=TEAL)
def divider(slide, t, w=12.6, c=TEAL): rect(slide, 0.55, t, w, 0.03, fill=c)
def tag(slide, label):
    rect(slide, 11.0, 0.22, 2.1, 0.34, fill=TEAL)
    txt(slide, label, 11.0, 0.22, 2.1, 0.34, size=9, bold=True, align=PP_ALIGN.CENTER)
def num(slide, n, total=12):
    txt(slide, f"{n} / {total}", 12.2, 7.1, 1.0, 0.35, size=9, color=LGREY, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, size=12, color=LGREY, gap=0.42):
    for i, item in enumerate(items):
        txt(slide, f"▸  {item}", l, t + i*gap, w, gap, size=size, color=color)

def kpi_box(slide, x, y, val, label, w=2.9, h=1.9, val_color=GOLD):
    rect(slide, x, y, w, h, fill=DGREY)
    rect(slide, x, y, w, 0.04, fill=TEAL)
    txt(slide, val,   x, y+0.25, w, 0.85, size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y+1.1, w, 0.7,  size=11, color=LGREY, align=PP_ALIGN.CENTER)

def card(slide, x, y, title, body, w=5.8, h=2.1):
    rect(slide, x, y, w, h, fill=DGREY)
    rect(slide, x, y, w, 0.42, fill=BLUE)
    txt(slide, title, x+0.15, y+0.04, w-0.2, 0.36, size=11, bold=True)
    txt(slide, body,  x+0.15, y+0.55, w-0.2, h-0.65, size=11, color=LGREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
rect(sl, 9.2, 0, 4.13, 7.5, fill=DGREY)
rect(sl, 11.5, 0, 1.83, 7.5, fill=TEAL)

# Logo block
rect(sl, 0.55, 0.5, 3.2, 0.7, fill=TEAL)
txt(sl, "innOFocus", 0.55, 0.5, 3.2, 0.7, size=22, bold=True, align=PP_ALIGN.CENTER)

txt(sl, "CORPORATE FINANCING PITCH", 0.55, 1.65, 8.0, 0.45,
    size=12, bold=False, color=TEAL, italic=True)
txt(sl, "Innofocus Photonics Technology", 0.55, 2.1, 8.5, 0.9,
    size=38, bold=True, color=WHITE)
txt(sl, "Innovator of Smart 3D Laser Nanoprinting &\nCharacterisation Equipment",
    0.55, 3.15, 8.2, 0.85, size=16, color=LGREY)

divider(sl, 4.15, w=4.0, c=GOLD)

txt(sl, "Melbourne, Australia  ·  2026", 0.55, 4.45, 6.0, 0.4, size=13, color=LGREY)
txt(sl, "CONFIDENTIAL — FOR DISCUSSION PURPOSES ONLY",
    0.55, 7.1, 9.5, 0.35, size=9, italic=True, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "OVERVIEW"); num(sl, 2)

txt(sl, "Executive Summary", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

boxes = [
    ("🏢  The Company",
     ["Innofocus Photonics Technology Pty Ltd",
      "Headquartered in Melbourne, Australia",
      "1,000 m² Nano Manufacturing Plant",
      "InnovationAus 2025: Translation Hero Award\n      & People's Choice Award"]),
    ("🔬  The Technology",
     ["World's only commercial 3D refractive index\n      imaging system (HoloView 3DRI)",
      "Sub-100 nm resolution",
      "400× faster than competing systems",
      "Multi-material capability"]),
    ("💰  The Opportunity",
     ["Serving ICT, quantum tech & medical devices",
      "Global nanofabrication equipment market\n      growing at double-digit CAGR",
      "Unique IP position — no direct competitor",
      "Strong academic & industry partnerships"]),
]
for i, (title, pts) in enumerate(boxes):
    x = 0.42 + i * 4.25
    rect(sl, x, 1.45, 3.95, 5.3, fill=DGREY)
    rect(sl, x, 1.45, 3.95, 0.50, fill=BLUE)
    txt(sl, title, x+0.15, 1.47, 3.65, 0.46, size=12, bold=True)
    for j, pt in enumerate(pts):
        txt(sl, f"• {pt}", x+0.2, 2.15+j*0.75, 3.55, 0.7, size=11, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "PROBLEM"); num(sl, 3)

txt(sl, "The Problem", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

txt(sl, "Current nanofabrication tools leave critical gaps:", 0.55, 1.2, 10.0, 0.4,
    size=14, bold=True, color=GOLD)

problems = [
    ("Slow throughput",
     "Traditional two-photon lithography and e-beam systems are too slow for\nindustrial-scale production — holding back commercialisation of quantum\nphotonics and advanced ICT devices."),
    ("Poor resolution & precision",
     "Most commercial systems cannot reliably achieve sub-100 nm features\nrequired for next-generation photonic integrated circuits and fibre Bragg\ngratings."),
    ("No in-situ characterisation",
     "Fabrication and measurement are separate steps, causing costly re-work\ncycles. There is no commercially available system offering real-time 3D\nrefractive index imaging during manufacturing."),
    ("Limited material versatility",
     "Existing tools are often constrained to specific material classes,\npreventing manufacturers from working across soft, hard, biological,\nand semiconductor materials on a single platform."),
]
for i, (title, body) in enumerate(problems):
    r, c = divmod(i, 2)
    x = 0.42 + c * 6.35
    y = 1.8  + r * 2.55
    rect(sl, x, y, 6.1, 2.3, fill=DGREY)
    rect(sl, x, y, 6.1, 0.04, fill=RED)
    txt(sl, f"⚠  {title}", x+0.2, y+0.12, 5.7, 0.40, size=12, bold=True, color=GOLD)
    txt(sl, body, x+0.2, y+0.6, 5.7, 1.55, size=11, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Solution
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "SOLUTION"); num(sl, 4)

txt(sl, "Our Solution", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

txt(sl, "\"Innovator of smart 3D laser nanoprinting and characterisation equipment —\nachieving perfect nanomanufacturing results has never been easier.\"",
    0.55, 1.2, 12.3, 0.9, size=14, italic=True, color=TEAL)

advantages = [
    ("Sub-100 nm\nResolution", "Surpassing the diffraction\nlimit for ultra-precise features"),
    ("400× Faster\nSpeed",     "Industrial throughput\nvs conventional systems"),
    ("In-Situ 3D\nImaging",    "World-only real-time\n3D refractive index imaging"),
    ("Multi-Material\nCapability", "Hard to soft, biological\nto semiconductor"),
]
for i, (val, desc) in enumerate(advantages):
    x = 0.42 + i * 3.15
    rect(sl, x, 2.35, 2.9, 2.1, fill=DGREY)
    rect(sl, x, 2.35, 2.9, 0.04, fill=TEAL)
    txt(sl, val,  x, 2.5,  2.9, 0.8, size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(sl, desc, x, 3.35, 2.9, 0.9, size=11, color=LGREY, align=PP_ALIGN.CENTER)

txt(sl, "Product Portfolio:", 0.55, 4.7, 4.0, 0.38, size=13, bold=True, color=GOLD)
products = [
    "nanoLAB Series — research-grade 3D laser nanoprinters",
    "nanoFACTORY Series — industrial nanomanufacturing platform",
    "nanoFACTORY i-QPC — quantum photonic chip fabrication",
    "nanoFACTORY r-FBG — automated fibre Bragg grating production (up to km length)",
    "HoloView 3DRI Benchtop — world's only commercial 3D refractive index imaging system",
    "Ultrafast FBGs (UFBGs) — femtosecond laser fabricated, reflectivity up to 99.99%",
]
bullets(sl, products, 0.55, 5.15, 12.3, size=11, gap=0.37)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Market Opportunity
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "MARKET"); num(sl, 5)

txt(sl, "Market Opportunity", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

markets = [
    ("TAM", "$42B+",  "Global Nanofabrication &\nPhotonics Equipment Market"),
    ("SAM", "$8B",    "Laser Nanofabrication,\nFBG & RI Imaging Systems"),
    ("SOM", "$250M+", "Near-Term Addressable\nRevenue Opportunity"),
]
for i, (lbl, val, desc) in enumerate(markets):
    y = 1.45 + i * 1.75
    w = 12.5 - i * 1.2
    rect(sl, 0.42, y, w, 1.5, fill=DGREY)
    rect(sl, 0.42, y, 1.75, 1.5, fill=TEAL if i==0 else BLUE if i==1 else RGBColor(0x00,0x7A,0x9E))
    txt(sl, lbl, 0.42, y+0.15, 1.75, 0.5, size=22, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, val, 2.35, y+0.05, 3.0, 0.65, size=30, bold=True, color=GOLD)
    txt(sl, desc, 2.35, y+0.75, w-2.1, 0.6, size=12, color=LGREY)

txt(sl, "Key Growth Drivers:", 0.55, 6.6, 5.0, 0.38, size=12, bold=True, color=GOLD)
drivers = ["Quantum computing & photonics commercialisation", "5G/6G advanced ICT infrastructure demand", "Medical device miniaturisation & biosensing growth"]
txt(sl, "  ·  ".join(drivers), 0.55, 6.95, 12.3, 0.38, size=11, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Products in Detail
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "PRODUCTS"); num(sl, 6)

txt(sl, "Product Portfolio", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

prods = [
    ("nanoLAB Series",
     "Research-grade intelligent 3D laser nanoprinting. Available with integrated HoloView 3DRI (H3D) for world-first in-situ refractive index imaging during fabrication. Target: universities, research institutes."),
    ("nanoFACTORY Series",
     "Industrial-scale nanomanufacturing platform. Variants include i-QPC for quantum photonic chip production and r-FBG for fully automated, AI-assisted fibre Bragg grating fabrication up to km length."),
    ("HoloView 3DRI Benchtop",
     "World's only commercially available equipment for high-resolution in-situ 3D refractive index distribution measurement. Enables non-destructive characterisation of photonic devices and biomaterials."),
    ("Ultrafast FBGs & Self-cooling Film",
     "Customisable femtosecond-laser-written fibre Bragg gratings (reflectivity up to 99.99%, high-temperature stable). Plus nanostructured self-cooling film using passive radiative cooling technology."),
]
for i, (title, body) in enumerate(prods):
    r, c = divmod(i, 2)
    x = 0.42 + c * 6.42
    y = 1.45 + r * 2.75
    rect(sl, x, y, 6.15, 2.5, fill=DGREY)
    rect(sl, x, y, 6.15, 0.44, fill=BLUE)
    txt(sl, title, x+0.18, y+0.05, 5.8, 0.38, size=12, bold=True)
    txt(sl, body,  x+0.18, y+0.55, 5.8, 1.78, size=11, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Traction & Milestones
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "TRACTION"); num(sl, 7)

txt(sl, "Traction & Milestones", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

kpis = [
    ("1,000 m²",  "Nano Manufacturing\nPlant (NMP)"),
    ("Sub-100 nm","Resolution\nCapability"),
    ("400×",      "Speed vs\nConventional Systems"),
    ("99.99%",    "Max FBG\nReflectivity"),
]
for i, (val, lbl) in enumerate(kpis):
    kpi_box(sl, 0.42 + i*3.18, 1.4, val, lbl)

txt(sl, "Key Milestones:", 0.55, 3.6, 4.0, 0.38, size=13, bold=True, color=GOLD)
milestones = [
    "Company founded — spin-out from RMIT University nanophotonics research group",
    "nanoLAB Series launched — first commercial intelligent 3D laser nanofabrication system",
    "HoloView 3DRI released — world's only commercial 3D refractive index imaging system",
    "nanoFACTORY i-QPC & r-FBG introduced — targeting quantum photonics and fibre sensing markets",
    "Nano Manufacturing Plant (NMP) operational — 1,000 m² facility in Melbourne",
    "InnovationAus 2025 Awards — Translation Hero Award & People's Choice Award",
    "Partnerships with Coherent, Olympus, RMIT, Swinburne, UCLA, Peking University",
]
bullets(sl, milestones, 0.55, 4.1, 12.3, size=11, gap=0.40)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Financial Projections
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "FINANCIALS"); num(sl, 8)

txt(sl, "Financial Projections", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)
txt(sl, "5-Year Forecast  (AUD millions)", 0.55, 1.15, 8.0, 0.38, size=12, italic=True, color=LGREY)

headers  = ["Metric", "FY1", "FY2", "FY3", "FY4", "FY5"]
col_w    = [3.2, 1.7, 1.7, 1.7, 1.7, 1.7]
rows_data = [
    ("Revenue",              "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("Gross Profit",         "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("EBITDA",               "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("Net Profit / (Loss)",  "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("Gross Margin",         "[X]%",  "[X]%",  "[X]%",  "[X]%",  "[X]%"),
]
x0, y0 = 0.35, 1.6
for j, (h, cw) in enumerate(zip(headers, col_w)):
    x = x0 + sum(col_w[:j])
    rect(sl, x, y0, cw, 0.44, fill=TEAL if j==0 else BLUE)
    txt(sl, h, x+0.1, y0+0.03, cw-0.1, 0.38, size=11, bold=True,
        align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

for r, row in enumerate(rows_data):
    y = y0 + 0.44 + r * 0.48
    fill = DGREY if r%2==0 else RGBColor(0x0E, 0x22, 0x35)
    for j, (cell, cw) in enumerate(zip(row, col_w)):
        x = x0 + sum(col_w[:j])
        rect(sl, x, y, cw, 0.46, fill=fill)
        c = GOLD if j>0 else LGREY
        txt(sl, cell, x+0.1, y+0.03, cw-0.1, 0.38, size=11, color=c,
            align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

txt(sl, "Revenue streams: Equipment sales (nanoLAB / nanoFACTORY / HoloView), consumables (UFBGs, films), service contracts, licensing.",
    0.35, 4.85, 12.5, 0.45, size=11, color=LGREY)
txt(sl, "* Projections require input from management. Placeholders shown — replace with audited / management accounts.",
    0.35, 7.1, 12.5, 0.35, size=9, italic=True, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Use of Funds
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "USE OF FUNDS"); num(sl, 9)

txt(sl, "Use of Funds", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)
txt(sl, "Total Raise:  $[X] Million AUD", 0.55, 1.2, 7.0, 0.5, size=20, bold=True, color=GOLD)

allocs = [
    ("[X]%", "$[X]M", "Manufacturing Scale-Up",
     "Expand Nano Manufacturing Plant capacity, new CNC and laser tooling for nanoFACTORY production"),
    ("[X]%", "$[X]M", "R&D — Next-Generation Products",
     "HoloView v2, nanoFACTORY i-QPC v2, AI-enhanced process control, new self-cooling materials"),
    ("[X]%", "$[X]M", "International Market Expansion",
     "Sales offices in USA, Europe, and Asia; channel partnerships; trade show presence"),
    ("[X]%", "$[X]M", "Talent & Operations",
     "Hire engineers, application scientists, and sales staff; strengthen IP and legal infrastructure"),
    ("[X]%", "$[X]M", "Working Capital & Reserves",
     "Inventory, customer credit terms, operational buffer for 18-month runway"),
]
for i, (pct, amt, cat, desc) in enumerate(allocs):
    y = 2.0 + i * 0.98
    rect(sl, 0.42, y, 1.25, 0.8, fill=TEAL)
    txt(sl, pct, 0.42, y+0.05, 1.25, 0.42, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(sl, amt, 0.42, y+0.45, 1.25, 0.3,  size=11, color=NAVY,  align=PP_ALIGN.CENTER)
    txt(sl, cat,  1.85, y+0.0,  4.5, 0.38, size=13, bold=True, color=WHITE)
    txt(sl, desc, 1.85, y+0.38, 11.0, 0.42, size=11, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Competitive Advantage
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "COMPETITION"); num(sl, 10)

txt(sl, "Competitive Landscape", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

comp_headers = ["Capability", "innOFocus", "Nanoscribe", "Heidelberg\nInstruments", "Standard\nFBG Systems"]
comp_rows = [
    ("Sub-100 nm resolution",              "✔", "✔", "✘", "✘"),
    ("400× faster throughput",             "✔", "✘", "✘", "✘"),
    ("In-situ 3D RI imaging (commercial)", "✔", "✘", "✘", "✘"),
    ("Multi-material capability",          "✔", "✔", "✘", "✘"),
    ("FBG (km-length, automated)",         "✔", "✘", "✘", "✔"),
    ("Quantum photonic chip ready",        "✔", "✘", "✔", "✘"),
    ("Australian IP & manufacturing",      "✔", "✘", "✘", "✘"),
]
col_w = [3.8, 1.95, 1.95, 1.95, 1.95]
x0, y0 = 0.42, 1.45

for j, (h, cw) in enumerate(zip(comp_headers, col_w)):
    x = x0 + sum(col_w[:j])
    fill = GOLD if j==1 else TEAL if j==0 else BLUE
    rect(sl, x, y0, cw, 0.5, fill=fill)
    fc = NAVY if j<=1 else WHITE
    txt(sl, h, x+0.1, y0+0.02, cw-0.1, 0.46, size=10, bold=True, color=fc,
        align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

for r, row in enumerate(comp_rows):
    y = y0 + 0.5 + r * 0.52
    fill = DGREY if r%2==0 else RGBColor(0x0E, 0x22, 0x35)
    for j, (cell, cw) in enumerate(zip(row, col_w)):
        x = x0 + sum(col_w[:j])
        bg2 = RGBColor(0x0A, 0x30, 0x50) if j==1 else fill
        rect(sl, x, y, cw, 0.5, fill=bg2)
        if cell == "✔":
            c = GOLD if j==1 else GREEN
        elif cell == "✘":
            c = RED
        else:
            c = LGREY
        txt(sl, cell, x+0.1, y+0.04, cw-0.1, 0.42, size=13, bold=(j==1),
            color=c, align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Team & Partners
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); bar(sl); tag(sl, "TEAM"); num(sl, 11)

txt(sl, "Leadership & Partners", 0.55, 0.3, 9.0, 0.6, size=28, bold=True)
divider(sl, 1.05)

team = [
    ("[Founder / CEO]",  "Chief Executive Officer",
     "Deep expertise in photonics commercialisation; led company from RMIT spin-out to international sales"),
    ("[CTO]",            "Chief Technology Officer",
     "PhD in nanophotonics; architect of HoloView 3DRI and nanoFACTORY platform; multiple patents held"),
    ("[CFO]",            "Chief Financial Officer",
     "Background in deep-tech and hardware company financing; experienced in equity and debt raising"),
    ("[COO]",            "Chief Operating Officer",
     "Managed NMP scale-up; expertise in precision manufacturing and supply chain for photonics hardware"),
]
for i, (name, title, bio) in enumerate(team):
    r, c = divmod(i, 2)
    x = 0.42 + c*6.42; y = 1.42 + r*2.35
    rect(sl, x, y, 6.15, 2.1, fill=DGREY)
    rect(sl, x+0.15, y+0.25, 0.85, 0.85, fill=TEAL)
    txt(sl, "👤", x+0.17, y+0.25, 0.85, 0.85, size=22, align=PP_ALIGN.CENTER)
    txt(sl, name,  x+1.2, y+0.18, 4.7, 0.42, size=13, bold=True)
    txt(sl, title, x+1.2, y+0.6,  4.7, 0.38, size=11, color=TEAL)
    txt(sl, bio,   x+0.15, y+1.25, 5.8, 0.7,  size=10, color=LGREY)

txt(sl, "Strategic Partners & Collaborators:", 0.55, 6.42, 5.5, 0.38, size=12, bold=True, color=GOLD)
partners = ["RMIT University", "Swinburne University", "UCLA", "Peking University",
            "Coherent", "Olympus"]
txt(sl, "  ·  ".join(partners), 0.55, 6.8, 12.3, 0.4, size=12, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Investment Ask
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
rect(sl, 9.0, 0, 4.33, 7.5, fill=DGREY)
rect(sl, 11.2, 0, 2.13, 7.5, fill=TEAL)
bar(sl); num(sl, 12)

txt(sl, "THE ASK", 0.55, 1.1, 6.0, 0.45, size=13, bold=True, color=TEAL, italic=True)
txt(sl, "Partner with innOFocus", 0.55, 1.6, 8.3, 0.85, size=32, bold=True, color=WHITE)
txt(sl, "to lead the global nanofabrication revolution.", 0.55, 2.5, 8.3, 0.5, size=18, color=LGREY)
divider(sl, 3.2, w=5.0, c=GOLD)

ask_items = [
    ("Raise Amount:",    "$[X] Million AUD"),
    ("Instrument:",      "[Equity / Convertible Note / Term Loan]"),
    ("Pre-Money Valuation:", "$[X]M AUD"),
    ("Target Close:",    "[Q? 2026]"),
    ("Lead Investor:",   "[Open / Anchor TBC]"),
    ("Use of Proceeds:", "Manufacturing scale-up, R&D, global expansion"),
]
for i, (label, value) in enumerate(ask_items):
    y = 3.5 + i * 0.60
    txt(sl, label, 0.55, y, 3.6, 0.5, size=11, color=LGREY)
    txt(sl, value, 4.2,  y, 4.6, 0.5, size=12, bold=True)

txt(sl, "Contact:", 0.55, 7.0, 2.0, 0.35, size=11, bold=True, color=GOLD)
txt(sl, "service@innofocus.com.au  ·  +61 3 9077 8119  ·  innofocus.com.au",
    0.55, 7.15, 9.0, 0.3, size=11, color=LGREY)

# ── Save ────────────────────────────────────────────────────────────────────
OUT = "/tmp/innofocus_pitch.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
