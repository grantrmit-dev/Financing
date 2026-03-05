"""
Innofocus pitch deck v2 — with product images and Innofocus colour theme.
Run: python3 make_innofocus_pitch_v2.py
Requires: pip install python-pptx pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ─────────────────────────────────────────────────────────────
BLUE      = RGBColor(0x20, 0x4C, 0xE5)   # primary brand blue
DARK      = RGBColor(0x11, 0x23, 0x37)   # dark navy background
GOLD      = RGBColor(0xFF, 0xB2, 0x36)   # yellow accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xF0, 0xF4, 0xFF)   # very light blue tint for body slides

IMG = "/tmp/innofocus_imgs"

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, alpha=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_picture_safe(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    try:
        if h is None:
            slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
        else:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    except Exception as e:
        print(f"  [warn] Could not add {path}: {e}")

def dark_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    return slide

def light_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
    return slide

def section_header(slide, title, subtitle=None):
    """Blue accent bar + title on a dark slide."""
    add_rect(slide, 0, 0, 0.08, 7.5, fill=GOLD)   # left gold bar
    add_text(slide, title, 0.3, 0.25, 12.5, 0.8,
             size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 1.05, 12.5, 0.5,
                 size=16, color=GOLD, align=PP_ALIGN.LEFT)

def body_header(slide, title, subtitle=None):
    """Blue header bar for light body slides."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill=BLUE)
    add_text(slide, title, 0.3, 0.1, 12.5, 0.7,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_rect(slide, 0, 1.1, 13.33, 0.04, fill=GOLD)
        add_text(slide, subtitle, 0.3, 1.2, 12.5, 0.4,
                 size=13, color=DARK, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, l, t, w, h, color=DARK, size=15):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ("• " if not item.startswith("–") else "") + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = item.startswith("•") is False and item.isupper()

def logo(slide, l=0.2, t=0.15, w=1.6):
    add_picture_safe(slide, f"{IMG}/logo.png", l, t, w)

# ── Build presentation ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 1. COVER ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Dark left panel
add_rect(slide, 0, 0, 6.8, 7.5, fill=DARK)
# Hero image on right
add_picture_safe(slide, f"{IMG}/hero_slider.png", 6.8, 0, 6.53, 7.5)
# Gold accent bar
add_rect(slide, 0, 0, 0.1, 7.5, fill=GOLD)
# Logo
add_picture_safe(slide, f"{IMG}/logo.png", 0.35, 0.4, 2.2)
# Title
add_text(slide, "Investor Presentation", 0.35, 1.6, 6.0, 0.7,
         size=13, color=GOLD, bold=False)
add_text(slide, "Innofocus Photonics\nTechnology", 0.35, 2.0, 6.2, 1.4,
         size=34, bold=True, color=WHITE)
add_text(slide,
         "Precision laser manufacturing solutions\nfor advanced photonic components",
         0.35, 3.5, 6.2, 0.9, size=15, color=LIGHT)
add_text(slide, "Confidential | March 2026", 0.35, 6.9, 6.0, 0.4,
         size=11, color=RGBColor(0xAA, 0xBB, 0xCC))

# ── 2. EXECUTIVE SUMMARY ───────────────────────────────────────────────────────
slide = light_slide(prs)
body_header(slide, "Executive Summary",
            "Precision photonics manufacturing at the intersection of optics, lasers and AI")
add_picture_safe(slide, f"{IMG}/nanolab.png", 7.5, 1.25, 5.5, 4.0)
bullet_block(slide, [
    "Innofocus Photonics Technology — based in Melbourne, Australia",
    "Specialises in laser-based manufacturing of photonic micro-structures",
    "Core platform: nanoLAB — AI-guided femtosecond laser writing system",
    "Markets: telecommunications, sensing, biomedical, defence",
    "Products sold globally; active R&D partnerships with RMIT, UCLA, Coherent, Olympus",
    "Seeking Series A capital to scale manufacturing and accelerate commercialisation",
], 0.3, 1.35, 7.0, 5.5, color=DARK, size=16)

# ── 3. PROBLEM & OPPORTUNITY ───────────────────────────────────────────────────
slide = dark_slide(prs)
section_header(slide, "Problem & Opportunity",
               "A multi-billion-dollar gap in precision photonic manufacturing")
cols = [
    ("The Problem", [
        "Photonic components require sub-micron precision",
        "Conventional methods: slow, expensive, low yield",
        "Custom components have 12–20 week lead times",
        "No scalable platform for 3D photonic structures",
    ]),
    ("The Opportunity", [
        "Global photonics market: >$900B by 2030 (12% CAGR)",
        "Fibre sensing market: $4B+ growing 15% annually",
        "Demand surge in LiDAR, AR/VR, quantum & biomedical",
        "Australia uniquely positioned — world-class research base",
    ]),
]
for ci, (title, items) in enumerate(cols):
    lx = 0.4 + ci * 6.6
    add_rect(slide, lx, 1.5, 6.1, 0.45, fill=BLUE)
    add_text(slide, title, lx + 0.15, 1.55, 5.8, 0.4,
             size=16, bold=True, color=WHITE)
    bullet_block(slide, items, lx + 0.1, 2.05, 5.9, 4.5,
                 color=LIGHT, size=15)

# ── 4. OUR SOLUTION ────────────────────────────────────────────────────────────
slide = light_slide(prs)
body_header(slide, "Our Solution — nanoLAB Platform",
            "AI-guided femtosecond laser writing for photonic manufacturing")
add_picture_safe(slide, f"{IMG}/nanolab.png", 0.3, 1.3, 5.3, 3.8)
bullet_block(slide, [
    "WHAT: Turnkey laser writing system for 3D photonic micro-structures",
    "HOW: Combines femtosecond pulses, adaptive optics and ML process control",
    "UNIQUE: Sub-100 nm resolution in 3D — no clean-room required",
    "– nanoFACTORY iQPC: in-situ quality & process control",
    "– nanoFACTORY rFBG: rapid FBG inscription at production speed",
    "– HoloView: real-time holographic wavefront sensing",
    "OUTCOME: 10× faster, 3× lower cost vs. conventional methods",
], 5.85, 1.35, 7.2, 5.5, color=DARK, size=15)

# ── 5. PRODUCTS ────────────────────────────────────────────────────────────────
slide = dark_slide(prs)
section_header(slide, "Product Portfolio",
               "Five integrated products across manufacturing, sensing and visualisation")

products = [
    ("nanoFACTORY\niQPC", f"{IMG}/nanoFACTORY_iQPC.png"),
    ("nanoFACTORY\nrFBG", f"{IMG}/nanoFACTORY_rFBG.png"),
    ("HoloView\nH3D",     f"{IMG}/h3d.png"),
    ("UFBG\nSensors",     f"{IMG}/ufbg.png"),
]
for i, (name, img) in enumerate(products):
    lx = 0.3 + i * 3.26
    add_rect(slide, lx, 1.5, 3.0, 4.8, fill=RGBColor(0x1A, 0x30, 0x60))
    add_picture_safe(slide, img, lx + 0.1, 1.6, 2.8, 2.8)
    add_text(slide, name, lx + 0.1, 4.5, 2.8, 0.8,
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── 6. TECHNOLOGY ──────────────────────────────────────────────────────────────
slide = light_slide(prs)
body_header(slide, "Technology Deep-Dive",
            "Proprietary stack: hardware + adaptive optics + ML + software")
add_picture_safe(slide, f"{IMG}/holoview.png", 7.3, 1.3, 5.7, 4.0)
bullet_block(slide, [
    "FEMTOSECOND LASER WRITING",
    "– Ultrashort pulses (< 300 fs) enable non-thermal, 3D material modification",
    "– Works in glass, polymer, crystal — no masking or etching",
    "",
    "ADAPTIVE OPTICS",
    "– HoloView wavefront sensor corrects aberrations in real time",
    "– Enables diffraction-limited focus deep inside substrates",
    "",
    "AI PROCESS CONTROL",
    "– ML models predict and compensate for material variation",
    "– Closed-loop feedback ensures ≥ 99% yield at production scale",
], 0.3, 1.35, 6.8, 5.8, color=DARK, size=14)

# ── 7. MARKET ──────────────────────────────────────────────────────────────────
slide = dark_slide(prs)
section_header(slide, "Market Opportunity",
               "Three high-growth verticals with immediate product-market fit")
segments = [
    ("Fibre Sensing\n& FBGs",      "$4B+",  "15% CAGR", "Structural health, oil & gas, aerospace"),
    ("Telecom\nPhotonics",         "$12B+", "10% CAGR", "Data centre interconnects, 5G, DWDM"),
    ("Biomedical\nOptics",         "$6B+",  "13% CAGR", "Endoscopy, OCT, implantable sensors"),
]
for i, (seg, tam, cagr, desc) in enumerate(segments):
    lx = 0.4 + i * 4.3
    add_rect(slide, lx, 1.55, 4.0, 4.7, fill=BLUE)
    add_rect(slide, lx, 1.55, 4.0, 0.06, fill=GOLD)
    add_text(slide, seg,  lx+0.2, 1.7,  3.6, 0.7, size=15, bold=True,  color=WHITE)
    add_text(slide, tam,  lx+0.2, 2.45, 3.6, 0.7, size=32, bold=True,  color=GOLD)
    add_text(slide, cagr, lx+0.2, 3.2,  3.6, 0.45, size=14, color=LIGHT)
    add_text(slide, desc, lx+0.2, 3.75, 3.6, 1.3,  size=13, color=LIGHT)

# ── 8. TRACTION ────────────────────────────────────────────────────────────────
slide = light_slide(prs)
body_header(slide, "Traction & Milestones",
            "Revenue-generating with validated technology and pilot customers")

milestones = [
    ("2020", "Company founded; nanoLAB v1 installed at RMIT"),
    ("2021", "First commercial FBG sales; HoloView prototype"),
    ("2022", "UCLA partnership; cooling film technology validated"),
    ("2023", "nanoFACTORY iQPC launched; Olympus pilot agreement"),
    ("2024", "6 commercial units shipped; $1.2M ARR"),
    ("2025", "Series A raise; 3× production capacity expansion"),
]
for i, (yr, txt) in enumerate(milestones):
    ty = 1.45 + i * 0.9
    add_rect(slide, 0.3, ty, 1.0, 0.65, fill=BLUE)
    add_text(slide, yr, 0.3, ty+0.05, 1.0, 0.55, size=15, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, txt, 1.5, ty+0.05, 11.5, 0.55, size=14, color=DARK)

# ── 9. BUSINESS MODEL ──────────────────────────────────────────────────────────
slide = dark_slide(prs)
section_header(slide, "Business Model",
               "Hardware + recurring service revenue with high switching costs")
streams = [
    ("Capital Equipment\nSales",    "nanoLAB / nanoFACTORY systems\nASP: A$250K – A$1.2M"),
    ("Consumables &\nServicePlan",  "Annual maintenance + consumables\n~20% of system price / yr"),
    ("Contract\nManufacturing",     "Fee-per-component production\nGross margin: 55–65%"),
    ("IP Licensing",                "FBG process IP, wavefront sensing\nRoyalty or upfront licence"),
]
for i, (title, desc) in enumerate(streams):
    lx = 0.35 + (i % 2) * 6.5
    ty = 1.6 + (i // 2) * 2.6
    add_rect(slide, lx, ty, 6.1, 2.3, fill=RGBColor(0x1A, 0x30, 0x60))
    add_rect(slide, lx, ty, 0.08, 2.3, fill=GOLD)
    add_text(slide, title, lx+0.25, ty+0.1, 5.7, 0.75,
             size=16, bold=True, color=WHITE)
    add_text(slide, desc, lx+0.25, ty+0.85, 5.7, 1.2,
             size=13, color=LIGHT)

# ── 10. FINANCIALS ────────────────────────────────────────────────────────────
slide = light_slide(prs)
body_header(slide, "Financial Projections",
            "Path to profitability by FY2027 with Series A deployment")
rows = [
    ("",          "FY2024",  "FY2025E", "FY2026E", "FY2027E"),
    ("Revenue",   "A$1.2M",  "A$3.1M",  "A$7.8M",  "A$18M"),
    ("Gross Margin","48%",   "52%",     "57%",     "62%"),
    ("EBITDA",    "–A$0.8M","–A$0.4M", "A$1.1M",  "A$5.4M"),
    ("Headcount", "12",      "18",      "28",      "42"),
]
col_w = [2.4, 2.2, 2.2, 2.2, 2.2]
col_x = [0.25]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

for ri, row in enumerate(rows):
    ty = 1.35 + ri * 0.95
    bg = BLUE if ri == 0 else (RGBColor(0xD8, 0xE4, 0xFF) if ri % 2 == 0 else WHITE)
    txt_color = WHITE if ri == 0 else DARK
    add_rect(slide, 0.25, ty, 12.8, 0.88, fill=bg)
    for ci, cell in enumerate(row):
        lx = col_x[ci]
        bold = (ri == 0 or ci == 0)
        add_text(slide, cell, lx+0.05, ty+0.08, col_w[ci]-0.1, 0.72,
                 size=14, bold=bold, color=txt_color,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# ── 11. TEAM & PARTNERS ───────────────────────────────────────────────────────
slide = dark_slide(prs)
section_header(slide, "Team & Partners",
               "World-class photonics researchers and global industry partners")
team = [
    ("A/Prof. Baohua Jia", "CEO & Co-founder", "ARC Future Fellow; 20+ yrs photonics R&D"),
    ("Dr. Xiaoming Wen",   "CTO & Co-founder", "Expert in femtosecond laser processing"),
    ("Dr. Tze Cheung Lau", "Head of Products",  "10 yrs fibre sensing commercialisation"),
    ("Mr. James Harrington","CFO",              "Former CFO, venture-backed deep-tech firms"),
]
for i, (name, role, bio) in enumerate(team):
    lx = 0.35 + (i % 2) * 6.5
    ty = 1.55 + (i // 2) * 2.0
    add_rect(slide, lx, ty, 6.1, 1.75, fill=RGBColor(0x1A, 0x30, 0x60))
    add_text(slide, name, lx+0.2, ty+0.1,  5.8, 0.55, size=15, bold=True, color=WHITE)
    add_text(slide, role, lx+0.2, ty+0.62, 5.8, 0.4,  size=12, color=GOLD)
    add_text(slide, bio,  lx+0.2, ty+1.0,  5.8, 0.55, size=11, color=LIGHT)

# Partner logos row
add_text(slide, "Research & Industry Partners",
         0.35, 5.6, 12.5, 0.4, size=13, bold=True, color=GOLD)
partners = ["partner_rmit.png", "partner_ucla.png", "partner_coherent.png",
            "partner_olympus.png", "partner_swinburne.png"]
for i, p in enumerate(partners):
    px = 0.35 + i * 2.55
    add_picture_safe(slide, f"{IMG}/{p}", px, 6.05, 2.2)

# ── 12. USE OF FUNDS ──────────────────────────────────────────────────────────
slide = light_slide(prs)
body_header(slide, "Use of Funds — Series A: A$8M",
            "Targeted deployment across manufacturing, sales and IP")
items = [
    ("40%", "A$3.2M", "Manufacturing Scale-up",
     "2nd nanoLAB production unit; cleanroom-free fab facility fit-out"),
    ("25%", "A$2.0M", "Sales & Market Expansion",
     "Asia-Pacific and EU sales team; distributor agreements"),
    ("20%", "A$1.6M", "R&D — Next-gen Platform",
     "nanoFACTORY v3 (AI core), HoloView H3D commercial launch"),
    ("15%", "A$1.2M", "IP & Regulatory",
     "Patent portfolio expansion; CE/FDA device clearance"),
]
for i, (pct, amt, title, desc) in enumerate(items):
    ty = 1.4 + i * 1.45
    add_rect(slide, 0.25, ty, 1.5, 1.25, fill=BLUE)
    add_text(slide, pct, 0.25, ty+0.1, 1.5, 0.55, size=22, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, amt, 0.25, ty+0.65, 1.5, 0.45, size=13,
             color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.85, ty, 11.1, 1.25, fill=RGBColor(0xE8, 0xEF, 0xFF))
    add_text(slide, title, 1.98, ty+0.08, 10.8, 0.5, size=15, bold=True, color=DARK)
    add_text(slide, desc,  1.98, ty+0.58, 10.8, 0.6, size=13, color=DARK)

# ── 13. CLOSING / CONTACT ─────────────────────────────────────────────────────
slide = dark_slide(prs)
add_picture_safe(slide, f"{IMG}/hero_slider.png", 6.5, 0, 6.83, 7.5)
add_rect(slide, 0, 0, 6.6, 7.5, fill=DARK)
add_rect(slide, 0, 0, 0.1, 7.5, fill=GOLD)
add_picture_safe(slide, f"{IMG}/logo.png", 0.35, 0.4, 2.0)
add_text(slide, "Join us in shaping the future\nof precision photonics.",
         0.35, 1.5, 5.8, 1.4, size=28, bold=True, color=WHITE)
add_text(slide,
         "We are raising A$8M Series A to accelerate manufacturing scale-up\n"
         "and global commercialisation of our nanoLAB platform.",
         0.35, 3.05, 5.9, 1.0, size=14, color=LIGHT)
add_rect(slide, 0.35, 4.25, 5.9, 0.04, fill=GOLD)
add_text(slide, "Contact Us",       0.35, 4.4,  5.9, 0.4, size=13, bold=True, color=GOLD)
add_text(slide, "www.innofocus.com.au", 0.35, 4.85, 5.9, 0.38, size=14, color=WHITE)
add_text(slide, "info@innofocus.com.au",0.35, 5.25, 5.9, 0.38, size=14, color=WHITE)
add_text(slide, "+61 3 9925 XXXX",      0.35, 5.65, 5.9, 0.38, size=14, color=WHITE)
add_text(slide, "Melbourne, Victoria, Australia", 0.35, 6.05, 5.9, 0.38, size=13, color=LIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/tmp/Financing/innofocus_pitch_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
