"""
Generate a professional corporate financing pitch deck using python-pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy background
BLUE   = RGBColor(0x1A, 0x6E, 0xC8)   # primary accent
GOLD   = RGBColor(0xF0, 0xA5, 0x00)   # highlight / numbers
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xC8, 0xD6, 0xE5)   # light text
DGREY  = RGBColor(0x2A, 0x3A, 0x4A)   # card background

W  = Inches(13.33)   # widescreen 16:9 width
H  = Inches(7.50)    # widescreen 16:9 height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ───────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line:
        shape.line.color.rgb = line
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_divider(slide, t, color=BLUE, width=12.0, thickness=0.03):
    add_rect(slide, 0.65, t, width, thickness, fill=color)

def navy_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)

def accent_bar(slide, height=0.55):
    """Left blue accent bar."""
    add_rect(slide, 0, 0, 0.22, 7.5, fill=BLUE)

def slide_number(slide, n, total=12):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 1.0, 0.35,
             size=9, color=LGREY, align=PP_ALIGN.RIGHT)

def section_tag(slide, label):
    """Top-right section tag pill."""
    add_rect(slide, 11.2, 0.25, 1.9, 0.32, fill=BLUE)
    add_text(slide, label, 11.2, 0.25, 1.9, 0.32,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def bullet(slide, items, l, t, w, size=14, color=WHITE, spacing=0.38):
    for i, item in enumerate(items):
        add_text(slide, f"▸  {item}", l, t + i * spacing, w, 0.36,
                 size=size, color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)

# Large blue triangle decoration
add_rect(sl, 8.5, 0, 4.83, 7.5, fill=DGREY)
add_rect(sl, 10.5, 0, 2.83, 7.5, fill=BLUE)

# Logo placeholder
add_rect(sl, 0.65, 0.5, 2.2, 0.7, fill=BLUE)
add_text(sl, "[COMPANY LOGO]", 0.65, 0.5, 2.2, 0.7, size=10, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "CORPORATE FINANCING", 0.65, 1.7, 7.0, 0.5,
         size=13, bold=False, color=GOLD, italic=True)
add_text(sl, "[Company Name]", 0.65, 2.2, 8.5, 1.0,
         size=42, bold=True, color=WHITE)
add_text(sl, "Investor Pitch Deck", 0.65, 3.3, 7.0, 0.5,
         size=20, color=LGREY)

add_divider(sl, 4.05, color=GOLD, width=3.5, thickness=0.04)

add_text(sl, "[City, Country]  ·  [Month Year]", 0.65, 4.3, 6.0, 0.4,
         size=13, color=LGREY)
add_text(sl, "CONFIDENTIAL — FOR DISCUSSION PURPOSES ONLY", 0.65, 7.1, 9.0, 0.35,
         size=9, italic=True, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "OVERVIEW")
slide_number(sl, 2)

add_text(sl, "Executive Summary", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

# 3 summary boxes
boxes = [
    ("🏢  The Company", ["[Company Name] is a [industry] company", "Founded [year] | HQ: [location]", "Serving [target market] globally"]),
    ("💰  The Ask",     ["Raising  [$ X M] in [debt / equity]", "Use of funds: growth, capex, M&A", "Target close: [Q? 20??]"]),
    ("📈  The Opportunity", ["Market size: $[X]B TAM", "[X]% CAGR projected through 20??", "Strong competitive moat"]),
]
for i, (title, pts) in enumerate(boxes):
    x = 0.5 + i * 4.25
    add_rect(sl, x, 1.5, 3.9, 4.8, fill=DGREY)
    add_rect(sl, x, 1.5, 3.9, 0.52, fill=BLUE)
    add_text(sl, title, x + 0.15, 1.52, 3.6, 0.48, size=12, bold=True, color=WHITE)
    for j, pt in enumerate(pts):
        add_text(sl, f"• {pt}", x + 0.2, 2.2 + j * 0.5, 3.5, 0.45, size=12, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem / Opportunity
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "PROBLEM")
slide_number(sl, 3)

add_text(sl, "The Problem & Opportunity", 0.55, 0.35, 10.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

add_text(sl, "Current market pain points:", 0.55, 1.3, 6.0, 0.4,
         size=14, bold=True, color=GOLD)
pain = [
    "[Pain point 1 — e.g. fragmented market, high costs, inefficiency]",
    "[Pain point 2 — e.g. lack of technology, regulatory burden]",
    "[Pain point 3 — e.g. underserved customer segment]",
]
bullet(sl, pain, 0.55, 1.8, 8.0, size=13)

add_rect(sl, 0.55, 3.25, 12.3, 0.02, fill=BLUE)

add_text(sl, "The Opportunity:", 0.55, 3.45, 6.0, 0.4,
         size=14, bold=True, color=GOLD)

# Stat boxes
stats = [("$[X]B", "Total Addressable\nMarket"), ("[X]%", "Projected Annual\nGrowth Rate"), ("[X]M+", "Potential\nCustomers")]
for i, (num, lbl) in enumerate(stats):
    x = 0.55 + i * 4.1
    add_rect(sl, x, 4.0, 3.7, 2.8, fill=DGREY)
    add_text(sl, num, x, 4.4, 3.7, 1.0, size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x, 5.5, 3.7, 0.9, size=12, color=LGREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Solution / Business Model
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "SOLUTION")
slide_number(sl, 4)

add_text(sl, "Our Solution & Business Model", 0.55, 0.35, 10.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

add_text(sl, "[One-sentence value proposition: We help [target customer] to [achieve outcome] by [unique approach].]",
         0.55, 1.3, 12.3, 0.9, size=15, italic=True, color=GOLD)

# 4-box model
model = [
    ("Product / Service", "[Describe core offering, technology or IP]"),
    ("Revenue Model",     "[Subscription / SaaS / transaction fee / licensing / etc.]"),
    ("Key Differentiators", "[What makes us defensible vs. competitors]"),
    ("Go-to-Market",      "[Sales channels, partnerships, distribution strategy]"),
]
for i, (title, body) in enumerate(model):
    r, c = divmod(i, 2)
    x = 0.55 + c * 6.3
    y = 2.5  + r * 2.3
    add_rect(sl, x, y, 5.9, 2.0, fill=DGREY)
    add_rect(sl, x, y, 5.9, 0.42, fill=BLUE)
    add_text(sl, title, x + 0.15, y + 0.03, 5.6, 0.38, size=12, bold=True, color=WHITE)
    add_text(sl, body,  x + 0.15, y + 0.55, 5.6, 1.3,  size=12, color=LGREY, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Market Size
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "MARKET")
slide_number(sl, 5)

add_text(sl, "Market Size", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

# Concentric circle visual (text-based)
circles = [
    ("TAM", "$[X]B", "Total Addressable Market — global market for [industry]"),
    ("SAM", "$[X]B", "Serviceable Available Market — segments we can reach"),
    ("SOM", "$[X]M", "Serviceable Obtainable Market — realistic near-term capture"),
]
for i, (lbl, val, desc) in enumerate(circles):
    y = 1.5 + i * 1.7
    w = 12 - i * 1.5
    add_rect(sl, 0.55, y, w, 1.45, fill=DGREY)
    add_rect(sl, 0.55, y, 1.8, 1.45, fill=BLUE)
    add_text(sl, lbl, 0.55, y + 0.1, 1.8, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, val, 2.55, y + 0.05, 2.5, 0.6, size=28, bold=True, color=GOLD)
    add_text(sl, desc, 2.55, y + 0.7, w - 2.2, 0.6, size=12, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Traction & KPIs
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "TRACTION")
slide_number(sl, 6)

add_text(sl, "Traction & Key Milestones", 0.55, 0.35, 10.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

kpis = [
    ("[$ X M]", "Revenue\n(LTM)"),
    ("[X]%",    "YoY Revenue\nGrowth"),
    ("[X]k",    "Active\nCustomers"),
    ("[X]%",    "Gross\nMargin"),
]
for i, (val, lbl) in enumerate(kpis):
    x = 0.55 + i * 3.1
    add_rect(sl, x, 1.4, 2.8, 2.0, fill=DGREY)
    add_rect(sl, x, 1.4, 2.8, 0.04, fill=GOLD)
    add_text(sl, val, x, 1.9, 2.8, 0.9, size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x, 2.85, 2.8, 0.7, size=12, color=LGREY, align=PP_ALIGN.CENTER)

add_text(sl, "Key Milestones:", 0.55, 3.75, 5.0, 0.4, size=14, bold=True, color=GOLD)
milestones = [
    "[20??] — Founded / product launched",
    "[20??] — First $[X]M revenue milestone",
    "[20??] — Expansion to [market/region]",
    "[20??] — Strategic partnership with [partner]",
    "[20?? target] — This financing round closes",
]
bullet(sl, milestones, 0.55, 4.25, 12.0, size=12, spacing=0.40)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Financial Projections
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "FINANCIALS")
slide_number(sl, 7)

add_text(sl, "Financial Projections", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

add_text(sl, "5-Year Forecast  (USD millions)", 0.55, 1.25, 8.0, 0.4,
         size=12, italic=True, color=LGREY)

# Table header
headers = ["Metric", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]
col_w = [3.0, 1.8, 1.8, 1.8, 1.8, 1.8]
rows_data = [
    ("Revenue",         "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("Gross Profit",    "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("EBITDA",          "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("Net Income",      "$[X]M", "$[X]M", "$[X]M", "$[X]M", "$[X]M"),
    ("Gross Margin %",  "[X]%",  "[X]%",  "[X]%",  "[X]%",  "[X]%"),
]

x_start = 0.35
y_header = 1.7
for j, (h, cw) in enumerate(zip(headers, col_w)):
    x = x_start + sum(col_w[:j])
    add_rect(sl, x, y_header, cw, 0.42, fill=BLUE)
    add_text(sl, h, x+0.1, y_header+0.02, cw-0.1, 0.38,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

for r, row in enumerate(rows_data):
    y = y_header + 0.42 + r * 0.44
    fill = DGREY if r % 2 == 0 else RGBColor(0x1A, 0x28, 0x38)
    for j, (cell, cw) in enumerate(zip(row, col_w)):
        x = x_start + sum(col_w[:j])
        add_rect(sl, x, y, cw, 0.42, fill=fill)
        c = GOLD if j > 0 else LGREY
        add_text(sl, cell, x+0.1, y+0.02, cw-0.1, 0.38,
                 size=11, color=c, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

add_text(sl, "* Projections are forward-looking estimates. Actual results may vary.",
         0.35, 6.95, 12.0, 0.35, size=9, italic=True, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Use of Funds
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "USE OF FUNDS")
slide_number(sl, 8)

add_text(sl, "Use of Funds", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

add_text(sl, f"Total Raise: $[X] Million", 0.55, 1.3, 6.0, 0.5,
         size=20, bold=True, color=GOLD)

allocations = [
    ("[X]%", "[$ X M]", "Product & Technology", "R&D, engineering headcount, platform development"),
    ("[X]%", "[$ X M]", "Sales & Marketing",    "Customer acquisition, brand building, partnerships"),
    ("[X]%", "[$ X M]", "Operations & CapEx",   "Infrastructure, equipment, facilities"),
    ("[X]%", "[$ X M]", "Working Capital",       "Inventory, receivables, operational buffer"),
    ("[X]%", "[$ X M]", "General & Admin",       "Corporate overhead, legal, compliance"),
]
for i, (pct, amt, cat, desc) in enumerate(allocations):
    y = 2.0 + i * 0.98
    bar_w = 1.2
    add_rect(sl, 0.55, y, bar_w, 0.75, fill=BLUE)
    add_text(sl, pct, 0.55, y+0.05, bar_w, 0.35, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, amt, 0.55, y+0.4,  bar_w, 0.3,  size=11, color=LGREY, align=PP_ALIGN.CENTER)
    add_text(sl, cat,  2.0, y+0.0,  5.0, 0.38, size=13, bold=True, color=WHITE)
    add_text(sl, desc, 2.0, y+0.38, 10.8, 0.38, size=11, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Competitive Landscape
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "COMPETITION")
slide_number(sl, 9)

add_text(sl, "Competitive Landscape", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

# Comparison table
comp_headers = ["Feature", "[Company]", "Competitor A", "Competitor B", "Competitor C"]
comp_rows = [
    ("[Key differentiator 1]", "✔", "✔", "✘", "✘"),
    ("[Key differentiator 2]", "✔", "✘", "✔", "✘"),
    ("[Key differentiator 3]", "✔", "✘", "✘", "✔"),
    ("[Key differentiator 4]", "✔", "✘", "✘", "✘"),
    ("Price competitiveness",  "✔", "✘", "✔", "✔"),
]
col_widths = [4.2, 2.0, 2.0, 2.0, 2.0]
x0 = 0.55
y0 = 1.5

for j, (h, cw) in enumerate(zip(comp_headers, col_widths)):
    x = x0 + sum(col_widths[:j])
    fill = GOLD if j == 1 else BLUE
    add_rect(sl, x, y0, cw, 0.45, fill=fill)
    fc = NAVY if j == 1 else WHITE
    add_text(sl, h, x+0.1, y0+0.02, cw-0.1, 0.4,
             size=11, bold=True, color=fc, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

for r, row in enumerate(comp_rows):
    y = y0 + 0.45 + r * 0.48
    fill = DGREY if r % 2 == 0 else RGBColor(0x1A, 0x28, 0x38)
    for j, (cell, cw) in enumerate(zip(row, col_widths)):
        x = x0 + sum(col_widths[:j])
        bg = RGBColor(0x0D, 0x30, 0x50) if j == 1 else fill
        add_rect(sl, x, y, cw, 0.46, fill=bg)
        col = GOLD if (j == 1 and cell == "✔") else (RGBColor(0x4A, 0xDE, 0x80) if cell == "✔" else RGBColor(0xF8, 0x71, 0x71) if cell == "✘" else LGREY)
        add_text(sl, cell, x+0.1, y+0.04, cw-0.1, 0.38,
                 size=13, bold=(j==1), color=col, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Team
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "TEAM")
slide_number(sl, 10)

add_text(sl, "Leadership Team", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

team = [
    ("[Name]", "Chief Executive Officer", "[Background: former CEO of X, X years in industry, MBA from Y]"),
    ("[Name]", "Chief Financial Officer", "[Background: ex-investment banker, CPA, prior CFO experience]"),
    ("[Name]", "Chief Technology Officer", "[Background: engineering PhD, previously VP Eng at X, Y patents]"),
    ("[Name]", "Chief Revenue Officer", "[Background: built $Xm sales team at X, deep industry network]"),
]
advisors = ["[Advisor 1] — [Relevant expertise]", "[Advisor 2] — [Relevant expertise]", "[Advisor 3] — [Relevant expertise]"]

for i, (name, title, bio) in enumerate(team):
    r, c = divmod(i, 2)
    x = 0.55 + c * 6.3
    y = 1.5  + r * 2.5
    add_rect(sl, x, y, 5.9, 2.1, fill=DGREY)
    # Avatar circle placeholder
    add_rect(sl, x + 0.2, y + 0.3, 0.9, 0.9, fill=BLUE)
    add_text(sl, "👤", x + 0.22, y + 0.3, 0.9, 0.9, size=24, align=PP_ALIGN.CENTER)
    add_text(sl, name,  x + 1.3, y + 0.2,  4.4, 0.45, size=14, bold=True, color=WHITE)
    add_text(sl, title, x + 1.3, y + 0.65, 4.4, 0.38, size=11, color=GOLD)
    add_text(sl, bio,   x + 0.2, y + 1.3,  5.5, 0.65, size=10, color=LGREY, wrap=True)

add_text(sl, "Advisory Board:", 0.55, 6.55, 4.0, 0.35, size=12, bold=True, color=GOLD)
add_text(sl, "  ·  ".join(advisors), 0.55, 6.9, 12.3, 0.4, size=10, color=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Risks & Mitigations
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
accent_bar(sl)
section_tag(sl, "RISKS")
slide_number(sl, 11)

add_text(sl, "Key Risks & Mitigations", 0.55, 0.35, 9.0, 0.6,
         size=28, bold=True, color=WHITE)
add_divider(sl, 1.1)

risks = [
    ("Market Risk",      "[Description of market/demand risk]",      "[Mitigation strategy: diversification, contracts, etc.]"),
    ("Execution Risk",   "[Operational or delivery risk]",           "[Mitigation: experienced team, phased milestones]"),
    ("Regulatory Risk",  "[Compliance or licensing risk]",           "[Mitigation: legal counsel, proactive engagement]"),
    ("Financial Risk",   "[Liquidity, FX, interest rate risk]",      "[Mitigation: hedging, covenant structure, reserves]"),
    ("Competitive Risk", "[New entrants or incumbent response]",     "[Mitigation: IP protection, customer lock-in, speed]"),
]

add_rect(sl, 0.55, 1.3,  4.5, 0.40, fill=BLUE)
add_rect(sl, 5.2,  1.3,  3.9, 0.40, fill=RGBColor(0x7C, 0x3A, 0x3A))
add_rect(sl, 9.25, 1.3,  3.7, 0.40, fill=RGBColor(0x1A, 0x5C, 0x3A))
add_text(sl, "Risk Category",  0.65, 1.3, 4.3, 0.40, size=11, bold=True, color=WHITE)
add_text(sl, "Risk Description", 5.3, 1.3, 3.7, 0.40, size=11, bold=True, color=WHITE)
add_text(sl, "Mitigation",     9.35, 1.3, 3.5, 0.40, size=11, bold=True, color=WHITE)

for r, (cat, risk, mit) in enumerate(risks):
    y = 1.70 + r * 0.92
    fill = DGREY if r % 2 == 0 else RGBColor(0x1A, 0x28, 0x38)
    add_rect(sl, 0.55, y, 4.5, 0.85, fill=fill)
    add_rect(sl, 5.2,  y, 3.9, 0.85, fill=fill)
    add_rect(sl, 9.25, y, 3.7, 0.85, fill=fill)
    add_text(sl, cat,  0.65, y+0.05, 4.3, 0.38, size=12, bold=True, color=GOLD)
    add_text(sl, risk,  5.3, y+0.05, 3.7, 0.75, size=11, color=LGREY, wrap=True)
    add_text(sl, mit,   9.35, y+0.05, 3.5, 0.75, size=11, color=RGBColor(0x4A, 0xDE, 0x80), wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Investment Ask & Close
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 9.0, 0, 4.33, 7.5, fill=DGREY)
add_rect(sl, 11.0, 0, 2.33, 7.5, fill=BLUE)
accent_bar(sl)
slide_number(sl, 12)

add_text(sl, "THE ASK", 0.55, 1.2, 6.0, 0.5,
         size=13, bold=True, color=GOLD, italic=True)
add_text(sl, "Invest in [Company Name]", 0.55, 1.7, 8.2, 1.0,
         size=34, bold=True, color=WHITE)
add_divider(sl, 2.85, color=GOLD, width=4.0, thickness=0.04)

ask_items = [
    ("Raise Amount:",    "$[X] Million"),
    ("Instrument:",      "[Equity / Convertible Note / Term Loan]"),
    ("Valuation:",       "$[Pre-money / Post-money] [X]M"),
    ("Target Close:",    "[Month Year]"),
    ("Lead Investor:",   "[Open / Anchor investor name]"),
]
for i, (label, value) in enumerate(ask_items):
    y = 3.1 + i * 0.72
    add_text(sl, label, 0.55, y, 3.5, 0.5, size=12, color=LGREY)
    add_text(sl, value, 4.0,  y, 5.0, 0.5, size=13, bold=True, color=WHITE)

add_text(sl, "Contact Us", 0.55, 6.55, 4.0, 0.4, size=12, bold=True, color=GOLD)
add_text(sl, "[Name]  ·  [email@company.com]  ·  [+XX XXX XXX XXXX]", 0.55, 6.9, 9.0, 0.4, size=11, color=LGREY)

# ── Save ───────────────────────────────────────────────────────────────────────
OUT = "/tmp/pitch_deck.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
